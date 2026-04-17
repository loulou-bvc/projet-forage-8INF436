"""
Script de génération du PowerPoint - Classification du Score de Crédit
Cours 8INF436 - Forage de données
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─── COULEURS ────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x1B, 0x2A, 0x4A)   # bleu marine foncé
GOLD       = RGBColor(0xF0, 0xA5, 0x00)   # or/ambre
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xFA)
MID_GRAY   = RGBColor(0x6C, 0x75, 0x8A)
GREEN_OK   = RGBColor(0x27, 0xAE, 0x60)
RED_BAD    = RGBColor(0xC0, 0x39, 0x2B)
ORANGE_MID = RGBColor(0xE6, 0x7E, 0x22)
DARK_BLUE  = RGBColor(0x0D, 0x1B, 0x35)

IMG_DIR = "/home/user/projet-forage-8INF436/pptx_images"

# ─── HELPERS ─────────────────────────────────────────────────────────────────

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, border_color=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_bullet_slide(slide, title_text, bullets, accent=GOLD):
    """Slide avec titre + liste de puces structurées."""
    set_bg(slide, WHITE)
    # bande de titre
    add_rect(slide, 0, 0, 13.33, 1.4, NAVY)
    # accent line
    add_rect(slide, 0, 1.4, 13.33, 0.06, accent)
    add_text(slide, title_text, 0.45, 0.22, 12.4, 1.1, 28, bold=True,
             color=WHITE, align=PP_ALIGN.LEFT)

    # numéro de section (petite barre dorée à gauche)
    add_rect(slide, 0.35, 1.7, 0.08, 5.5, accent)

    y = 1.75
    for item in bullets:
        if isinstance(item, dict) and item.get("type") == "header":
            add_text(slide, item["text"], 0.6, y, 12.4, 0.45,
                     13, bold=True, color=NAVY)
            y += 0.45
        else:
            label = item if isinstance(item, str) else item.get("text", "")
            # puce
            add_rect(slide, 0.6, y + 0.1, 0.12, 0.12, accent)
            add_text(slide, label, 0.85, y, 11.9, 0.38,
                     12.5, color=RGBColor(0x22, 0x2C, 0x45))
            y += 0.42
        if y > 6.8:
            break

def add_two_col_slide(slide, title_text, left_items, right_items,
                      left_title="", right_title=""):
    set_bg(slide, WHITE)
    add_rect(slide, 0, 0, 13.33, 1.4, NAVY)
    add_rect(slide, 0, 1.4, 13.33, 0.06, GOLD)
    add_text(slide, title_text, 0.45, 0.22, 12.4, 1.1, 28, bold=True,
             color=WHITE, align=PP_ALIGN.LEFT)

    for col, (col_items, col_title, x_start) in enumerate(
            [(left_items, left_title, 0.4), (right_items, right_title, 6.85)]):
        if col_title:
            add_rect(slide, x_start, 1.7, 5.9, 0.38, NAVY)
            add_text(slide, col_title, x_start + 0.15, 1.73, 5.6, 0.35,
                     12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_rect(slide, x_start, 2.15, 0.07, 4.8, GOLD)
        y = 2.2
        for item in col_items:
            add_rect(slide, x_start + 0.25, y + 0.11, 0.1, 0.1, GOLD)
            add_text(slide, item, x_start + 0.45, y, 5.3, 0.4,
                     12, color=RGBColor(0x22, 0x2C, 0x45))
            y += 0.42

def add_image_slide(slide, title_text, img_path, caption="",
                    img_l=1.0, img_t=1.55, img_w=11.33, img_h=4.8):
    set_bg(slide, WHITE)
    add_rect(slide, 0, 0, 13.33, 1.4, NAVY)
    add_rect(slide, 0, 1.4, 13.33, 0.06, GOLD)
    add_text(slide, title_text, 0.45, 0.22, 12.4, 1.1, 28, bold=True,
             color=WHITE, align=PP_ALIGN.LEFT)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(img_l), Inches(img_t),
                                 Inches(img_w), Inches(img_h))
    if caption:
        add_text(slide, caption, 0.5, 6.5, 12.33, 0.45, 10,
                 italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ─── CRÉATION DU FICHIER ─────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]   # blank

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – TITRE
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
set_bg(s1, DARK_BLUE)

# fond décoratif : bande dorée gauche
add_rect(s1, 0, 0, 0.35, 7.5, GOLD)
# bloc titre principal
add_rect(s1, 0.55, 1.4, 12.3, 0.07, GOLD)

add_text(s1, "Classification du Score de Crédit",
         0.65, 1.6, 12.0, 1.5, 38, bold=True, color=WHITE,
         align=PP_ALIGN.LEFT)
add_text(s1, "Forage de données appliqué au domaine bancaire",
         0.65, 3.2, 11.5, 0.6, 18, italic=True, color=GOLD,
         align=PP_ALIGN.LEFT)

add_rect(s1, 0.55, 3.9, 12.3, 0.04, RGBColor(0x3A, 0x4E, 0x72))

add_text(s1, "8INF436 – Forage de données",
         0.65, 4.1, 9, 0.45, 14, color=RGBColor(0xB0, 0xBE, 0xD4),
         align=PP_ALIGN.LEFT)
add_text(s1, "Frantxa Cabrejos  •  Loup-Djabril Le Bivic  •  Nathan Razafindratsima",
         0.65, 4.6, 12, 0.45, 13, color=RGBColor(0xB0, 0xBE, 0xD4),
         align=PP_ALIGN.LEFT)
add_text(s1, "Université du Québec à Chicoutimi  —  2025",
         0.65, 5.1, 10, 0.4, 12, color=MID_GRAY, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – CONTEXTE & PROBLÉMATIQUE
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
set_bg(s2, WHITE)
add_rect(s2, 0, 0, 13.33, 1.4, NAVY)
add_rect(s2, 0, 1.4, 13.33, 0.06, GOLD)
add_text(s2, "Contexte & Problématique", 0.45, 0.22, 12.4, 1.1, 28,
         bold=True, color=WHITE)

# Carte gauche - contexte
add_rect(s2, 0.4, 1.65, 5.9, 5.2, LIGHT_GRAY)
add_rect(s2, 0.4, 1.65, 5.9, 0.55, NAVY)
add_text(s2, "CONTEXTE", 0.5, 1.7, 5.7, 0.45, 13, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)

context_items = [
    "Les banques accordent des crédits à des milliers de clients",
    "Évaluer le risque financier est une étape critique",
    "Un mauvais score peut mener à des pertes importantes",
    "L'humain seul ne peut traiter 100 000 dossiers efficacement",
    "Le machine learning automatise et fiabilise ce processus",
]
y = 2.35
for item in context_items:
    add_rect(s2, 0.7, y + 0.1, 0.1, 0.1, GOLD)
    add_text(s2, item, 0.92, y, 5.1, 0.42, 12,
             color=RGBColor(0x22, 0x2C, 0x45))
    y += 0.47

# Carte droite - objectif
add_rect(s2, 6.9, 1.65, 5.9, 5.2, LIGHT_GRAY)
add_rect(s2, 6.9, 1.65, 5.9, 0.55, NAVY)
add_text(s2, "OBJECTIF DU PROJET", 7.0, 1.7, 5.7, 0.45, 13, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)

# badges de classes
for i, (label, color, desc) in enumerate([
    ("GOOD",     GREEN_OK,   "Faible risque · Client fiable"),
    ("POOR",     ORANGE_MID, "Risque modéré · Surveillance recommandée"),
    ("STANDARD", RED_BAD,    "Risque élevé · Crédit risqué"),
]):
    y_b = 2.35 + i * 1.35
    add_rect(s2, 7.1, y_b, 1.4, 0.48, color)
    add_text(s2, label, 7.1, y_b + 0.05, 1.4, 0.38, 12, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s2, desc, 8.65, y_b + 0.05, 3.9, 0.4, 11.5,
             color=RGBColor(0x22, 0x2C, 0x45))

add_text(s2, "→ Prédire automatiquement la classe de crédit d'un client bancaire",
         7.1, 5.7, 5.5, 0.55, 12, bold=True, color=NAVY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – LE JEU DE DONNÉES
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
add_image_slide(s3,
    "Le Jeu de Données",
    os.path.join(IMG_DIR, "img_01_cell13.png"),
    caption="Distribution des classes : Standard 53.2% · Poor 29.0% · Good 17.8%  —  Déséquilibre de classes ~3x",
    img_l=1.2, img_t=1.55, img_w=10.9, img_h=4.65)

# stat boxes en haut à droite sur la photo  → on les met sous le titre
add_rect(s3, 0.4, 1.58, 2.5, 0.82, NAVY)
add_text(s3, "100 000", 0.45, 1.6, 2.4, 0.42, 20, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
add_text(s3, "observations", 0.45, 2.0, 2.4, 0.35, 11,
         color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s3, 3.1, 1.58, 2.5, 0.82, NAVY)
add_text(s3, "28 variables", 3.15, 1.6, 2.4, 0.42, 20, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
add_text(s3, "features (numériques + catégorielles)", 3.15, 2.0, 2.4, 0.35, 10,
         color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s3, 5.8, 1.58, 2.5, 0.82, NAVY)
add_text(s3, "3 classes", 5.85, 1.6, 2.4, 0.42, 20, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
add_text(s3, "Good · Poor · Standard", 5.85, 2.0, 2.4, 0.35, 11,
         color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – QUALITÉ & NETTOYAGE DES DONNÉES
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
add_image_slide(s4,
    "Qualité & Nettoyage des Données",
    os.path.join(IMG_DIR, "img_00_cell11.png"),
    caption="Aucune colonne ne dépasse 60% de valeurs manquantes → toutes les colonnes sont conservées",
    img_l=0.6, img_t=2.6, img_w=12.1, img_h=3.9)

# puces explicatives entre titre et graphe
bullets_clean = [
    "Données volontairement bruitées (valeurs aberrantes, caractères parasites, types mixtes)",
    "Regex appliqué pour nettoyer les champs numériques corrompus (âge négatif, revenus avec '_')",
    "Valeurs aberrantes supprimées (âge < 18 ou > 100, revenus négatifs)",
]
y = 1.6
for b in bullets_clean:
    add_rect(s4, 0.55, y + 0.11, 0.12, 0.12, GOLD)
    add_text(s4, b, 0.8, y, 12.3, 0.4, 12, color=RGBColor(0x22, 0x2C, 0x45))
    y += 0.44

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – PRÉTRAITEMENT DES DONNÉES
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
set_bg(s5, WHITE)
add_rect(s5, 0, 0, 13.33, 1.4, NAVY)
add_rect(s5, 0, 1.4, 13.33, 0.06, GOLD)
add_text(s5, "Prétraitement des Données", 0.45, 0.22, 12.4, 1.1, 28,
         bold=True, color=WHITE)

steps = [
    ("1", "Encodage catégoriel",
     "LabelEncoder sur 5 variables catégorielles (profession, type de prêt, comportement de paiement…)"),
    ("2", "Imputation",
     "Remplacement des valeurs manquantes par la médiane — robuste aux valeurs extrêmes"),
    ("3", "Normalisation",
     "MinMaxScaler → toutes les valeurs ramenées entre 0 et 1, indispensable pour le SVM"),
    ("4", "Découpage stratifié",
     "80% entraînement / 20% test — la proportion de chaque classe est preservée dans les deux ensembles"),
]

for i, (num, title, desc) in enumerate(steps):
    col = i % 2
    row = i // 2
    x = 0.45 + col * 6.45
    y = 1.7  + row * 2.6

    add_rect(s5, x, y, 6.0, 2.3, LIGHT_GRAY)
    add_rect(s5, x, y, 0.55, 2.3, NAVY)
    add_text(s5, num, x, y + 0.6, 0.55, 0.9, 26, bold=True,
             color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s5, title, x + 0.65, y + 0.15, 5.15, 0.5, 14, bold=True,
             color=NAVY)
    add_text(s5, desc, x + 0.65, y + 0.7, 5.15, 1.4, 12,
             color=RGBColor(0x22, 0x2C, 0x45), wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – SMOTE : ÉQUILIBRAGE DES CLASSES
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
add_image_slide(s6,
    "Équilibrage des Classes — SMOTE",
    os.path.join(IMG_DIR, "img_02_cell26.png"),
    caption="SMOTE (Synthetic Minority Over-sampling Technique) — génère des exemples synthétiques dans l'espace des features pour équilibrer les classes d'entraînement",
    img_l=0.8, img_t=2.35, img_w=11.7, img_h=4.0)

add_text(s6,
    "Problème : la classe 'Standard' est 3× plus représentée que 'Good' → le modèle ignorerait les minorités",
    0.55, 1.58, 12.2, 0.42, 12.5, bold=True, color=RED_BAD)
add_text(s6,
    "Solution : SMOTE appliqué uniquement sur l'ensemble d'entraînement (80 000 → 127 617 exemples)",
    0.55, 2.0, 12.2, 0.38, 12.5, color=RGBColor(0x22, 0x2C, 0x45))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – RÉDUCTION DE DIMENSIONNALITÉ
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
set_bg(s7, WHITE)
add_rect(s7, 0, 0, 13.33, 1.4, NAVY)
add_rect(s7, 0, 1.4, 13.33, 0.06, GOLD)
add_text(s7, "Réduction de Dimensionnalité", 0.45, 0.22, 12.4, 1.1, 28,
         bold=True, color=WHITE)

# Colonne gauche – Feature selection
add_rect(s7, 0.35, 1.6, 6.1, 5.65, LIGHT_GRAY)
add_rect(s7, 0.35, 1.6, 6.1, 0.52, NAVY)
add_text(s7, "Sélection de variables", 0.45, 1.63, 5.9, 0.45, 13.5,
         bold=True, color=GOLD, align=PP_ALIGN.CENTER)

img_feat = os.path.join(IMG_DIR, "img_03_cell29.png")
if os.path.exists(img_feat):
    s7.shapes.add_picture(img_feat, Inches(0.4), Inches(2.25), Inches(6.0), Inches(2.8))

fs_bullets = [
    "Importances issues de la Forêt Aléatoire",
    "Top 8 features retenues sur 22",
    "→ Outstanding Debt, Interest Rate, Delay…",
    "Utilisé pour : RF et Gradient Boosting",
]
y = 5.15
for b in fs_bullets:
    add_rect(s7, 0.6, y + 0.1, 0.09, 0.09, GOLD)
    add_text(s7, b, 0.8, y, 5.4, 0.38, 11, color=RGBColor(0x22, 0x2C, 0x45))
    y += 0.4

# Colonne droite – PCA
add_rect(s7, 6.9, 1.6, 6.1, 5.65, LIGHT_GRAY)
add_rect(s7, 6.9, 1.6, 6.1, 0.52, NAVY)
add_text(s7, "ACP (PCA)", 7.0, 1.63, 5.9, 0.45, 13.5,
         bold=True, color=GOLD, align=PP_ALIGN.CENTER)

img_pca = os.path.join(IMG_DIR, "img_05_cell32.png")
if os.path.exists(img_pca):
    s7.shapes.add_picture(img_pca, Inches(6.95), Inches(2.25), Inches(6.0), Inches(2.8))

pca_bullets = [
    "StandardScaler requis avant PCA",
    "Critère : 95% de variance expliquée",
    "→ 22 dimensions → 20 composantes",
    "Utilisé pour : SVM (LinéaireSVC)",
]
y = 5.15
for b in pca_bullets:
    add_rect(s7, 7.1, y + 0.1, 0.09, 0.09, GOLD)
    add_text(s7, b, 7.3, y, 5.4, 0.38, 11, color=RGBColor(0x22, 0x2C, 0x45))
    y += 0.4

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – LES MODÈLES DE CLASSIFICATION
# ══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank_layout)
set_bg(s8, WHITE)
add_rect(s8, 0, 0, 13.33, 1.4, NAVY)
add_rect(s8, 0, 1.4, 13.33, 0.06, GOLD)
add_text(s8, "Les Modèles de Classification", 0.45, 0.22, 12.4, 1.1, 28,
         bold=True, color=WHITE)

models = [
    ("Forêt Aléatoire", GREEN_OK,
     ["Ensemble de 100 arbres de décision indépendants",
      "Chaque arbre vote → classe majoritaire retenue",
      "Naturellement robuste au surapprentissage",
      "Utilise les 8 features sélectionnées"]),
    ("SVM — LinearSVC", ORANGE_MID,
     ["Séparateur à vaste marge dans l'espace PCA",
      "LinearSVC choisi pour sa rapidité (O(n) vs O(n³))",
      "20 composantes principales en entrée",
      "Calibré pour fournir des probabilités"]),
    ("Gradient Boosting", RGBColor(0x27, 0x6F, 0xBF),
     ["Construction séquentielle d'arbres corrective",
      "Chaque arbre corrige les erreurs du précédent",
      "Profondeur 3 · taux d'apprentissage 0.1",
      "Utilise les 8 features sélectionnées"]),
]

for i, (name, color, bullets) in enumerate(models):
    x = 0.35 + i * 4.33
    add_rect(s8, x, 1.6, 4.0, 0.6, color)
    add_text(s8, name, x, 1.65, 4.0, 0.5, 14, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s8, x, 2.2, 4.0, 5.05, LIGHT_GRAY)
    y = 2.3
    for b in bullets:
        add_rect(s8, x + 0.22, y + 0.1, 0.1, 0.1, color)
        add_text(s8, b, x + 0.45, y, 3.35, 0.42, 12,
                 color=RGBColor(0x22, 0x2C, 0x45))
        y += 0.47

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – RÉSULTATS : COMPARAISON DES MODÈLES
# ══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank_layout)
add_image_slide(s9,
    "Résultats — Comparaison des Modèles",
    os.path.join(IMG_DIR, "img_06_cell44.png"),
    caption="Accuracy · Précision · Rappel · Score F1 (pondéré) — évaluation sur l'ensemble de test (20 000 exemples)",
    img_l=0.5, img_t=1.55, img_w=8.5, img_h=4.8)

# Tableau résumé côté droit
results = [
    ("Forêt Aléatoire",   "76.6%", "76.6%", GREEN_OK),
    ("Gradient Boosting", "68.5%", "68.7%", ORANGE_MID),
    ("SVM (LinearSVC)",   "53.2%", "52.9%", RED_BAD),
]
add_rect(s9, 9.2, 1.55, 3.85, 0.5, NAVY)
add_text(s9, "Modèle", 9.25, 1.58, 1.5, 0.4, 11, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
add_text(s9, "Acc.", 10.75, 1.58, 1.0, 0.4, 11, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
add_text(s9, "F1", 11.75, 1.58, 1.2, 0.4, 11, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)

for i, (model, acc, f1, color) in enumerate(results):
    y_r = 2.15 + i * 0.62
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(s9, 9.2, y_r, 3.85, 0.58, bg)
    add_rect(s9, 9.2, y_r, 0.18, 0.58, color)
    add_text(s9, model, 9.45, y_r + 0.1, 1.35, 0.38, 11,
             color=RGBColor(0x22, 0x2C, 0x45))
    add_text(s9, acc, 10.75, y_r + 0.1, 1.0, 0.38, 11, bold=True,
             color=color, align=PP_ALIGN.CENTER)
    add_text(s9, f1, 11.75, y_r + 0.1, 1.2, 0.38, 11, bold=True,
             color=color, align=PP_ALIGN.CENTER)

# Badge meilleur modèle
add_rect(s9, 9.2, 4.1, 3.85, 0.7, GREEN_OK)
add_text(s9, "MEILLEUR MODÈLE", 9.25, 4.13, 3.75, 0.3, 11, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s9, "Forêt Aléatoire  —  F1 = 76.6%", 9.25, 4.45, 3.75, 0.3, 12,
         bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s9,
    "Validation croisée (k=5) : RF atteint F1=82.4% · stabilité ±0.003",
    9.2, 5.0, 3.85, 0.5, 11, italic=True, color=MID_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – MATRICES DE CONFUSION
# ══════════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank_layout)
add_image_slide(s10,
    "Matrices de Confusion",
    os.path.join(IMG_DIR, "img_07_cell45.png"),
    caption="La diagonale = bonnes prédictions. La Forêt Aléatoire prédit correctement 8 213 'Standard', 4 528 'Poor' et 2 577 'Good'",
    img_l=0.3, img_t=2.15, img_w=12.7, img_h=4.45)

expl = [
    "La matrice de confusion montre pour chaque classe : combien d'exemples ont été bien ou mal classifiés",
    "La Forêt Aléatoire confond peu 'Good' et 'Poor' — erreurs principalement entre classes adjacentes",
    "Le SVM génère davantage d'erreurs croisées (moins bonne séparabilité dans l'espace PCA)",
]
y = 1.58
for b in expl:
    add_rect(s10, 0.5, y + 0.1, 0.12, 0.12, GOLD)
    add_text(s10, b, 0.75, y, 12.1, 0.42, 12, color=RGBColor(0x22, 0x2C, 0x45))
    y += 0.44

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – APPLICATION WEB STREAMLIT
# ══════════════════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(blank_layout)
set_bg(s11, WHITE)
add_rect(s11, 0, 0, 13.33, 1.4, NAVY)
add_rect(s11, 0, 1.4, 13.33, 0.06, GOLD)
add_text(s11, "Application Web — Streamlit", 0.45, 0.22, 12.4, 1.1, 28,
         bold=True, color=WHITE)

tabs = [
    ("Onglet 1", "Saisie manuelle",
     ["L'utilisateur entre les infos du client (âge, salaire, dettes…)",
      "Les 3 modèles prédisent chacun un score en temps réel",
      "Un vote majoritaire donne la décision finale",
      "Couleurs : vert (Good) · orange (Poor) · rouge (Standard)"]),
    ("Onglet 2", "Import CSV en masse",
     ["Téléversement d'un fichier CSV avec N clients",
      "Prédictions en batch pour tous les clients",
      "Export des résultats enrichis en CSV"]),
    ("Onglet 3", "Tableau de bord des métriques",
     ["Accuracy, Précision, Rappel, F1 et CV-F1",
      "Graphique comparatif des 3 modèles",
      "Identification automatique du meilleur modèle"]),
]

for i, (tag, title, bullets) in enumerate(tabs):
    x = 0.35 + i * 4.33
    add_rect(s11, x, 1.6, 4.0, 0.52, NAVY)
    add_text(s11, tag, x, 1.62, 1.3, 0.45, 10, color=GOLD,
             align=PP_ALIGN.CENTER)
    add_text(s11, title, x + 1.3, 1.62, 2.65, 0.45, 13, bold=True,
             color=WHITE, align=PP_ALIGN.LEFT)
    add_rect(s11, x, 2.12, 4.0, 5.15, LIGHT_GRAY)
    y = 2.22
    for b in bullets:
        add_rect(s11, x + 0.22, y + 0.1, 0.1, 0.1, GOLD)
        add_text(s11, b, x + 0.45, y, 3.35, 0.45, 12,
                 color=RGBColor(0x22, 0x2C, 0x45))
        y += 0.52

add_rect(s11, 0.35, 6.0, 12.63, 0.55, RGBColor(0xE8, 0xF0, 0xFE))
add_rect(s11, 0.35, 6.0, 0.22, 0.55, RGBColor(0x27, 0x6F, 0xBF))
add_text(s11,
    "Tous les modèles sont sérialisés (.pkl) — l'application charge le pipeline complet sans ré-entraînement",
    0.65, 6.05, 12.0, 0.45, 12, italic=True,
    color=RGBColor(0x22, 0x2C, 0x45))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(blank_layout)
set_bg(s12, DARK_BLUE)
add_rect(s12, 0, 0, 0.35, 7.5, GOLD)
add_rect(s12, 0.55, 1.3, 12.3, 0.06, GOLD)

add_text(s12, "Conclusion", 0.65, 0.25, 9, 1.0, 34, bold=True, color=WHITE)

# points bilan
bilan = [
    ("Pipeline complet",
     "Nettoyage → encodage → imputation → normalisation → SMOTE → réduction → classification"),
    ("3 modèles évalués",
     "Forêt Aléatoire (F1=76.6%) · Gradient Boosting (68.7%) · SVM (52.9%)"),
    ("Données déséquilibrées",
     "SMOTE a permis de traiter le déséquilibre sans biaiser l'évaluation sur le test"),
    ("Déploiement",
     "Interface Streamlit opérationnelle pour des prédictions individuelles et en masse"),
]

for i, (titre, texte) in enumerate(bilan):
    y = 1.6 + i * 1.25
    add_rect(s12, 0.6, y, 0.5, 0.5, GOLD)
    add_text(s12, str(i + 1), 0.6, y + 0.03, 0.5, 0.44, 18, bold=True,
             color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(s12, titre, 1.25, y + 0.0, 11.3, 0.38, 14, bold=True,
             color=GOLD)
    add_text(s12, texte, 1.25, y + 0.42, 11.3, 0.6, 12.5,
             color=RGBColor(0xC8, 0xD4, 0xE8))

# message final
add_rect(s12, 0.55, 6.55, 12.3, 0.65, RGBColor(0x0D, 0x28, 0x52))
add_text(s12,
    "Merci pour votre attention  —  Questions ?",
    0.6, 6.58, 12.2, 0.55, 18, bold=True, color=GOLD,
    align=PP_ALIGN.CENTER)

# ─── SAUVEGARDE ──────────────────────────────────────────────────────────────
out_path = "/home/user/projet-forage-8INF436/presentation_credit_score.pptx"
prs.save(out_path)
print(f"Fichier créé : {out_path}")
print(f"Nombre de slides : {len(prs.slides)}")
