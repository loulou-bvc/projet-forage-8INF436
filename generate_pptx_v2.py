"""
Présentation PowerPoint v2 — Design 2026
Classification du Score de Crédit · 8INF436
─────────────────────────────────────────
Principes appliqués (recherches 2026) :
  • Dark mode + accents dorés (tendance n°1 2026)
  • Bento-grid layout (boîtes modulaires)
  • Règle 10-20-30 (Guy Kawasaki) → ≤12 slides, police ≥18pt
  • Une seule idée centrale par slide
  • Les titres = le message, pas juste un label
  • Impact d'abord : chiffres clés en grand avant l'explication
  • Réduction maximale du texte → visuels > bullets
  • Hiérarchie visuelle claire
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

IMG = "/home/user/projet-forage-8INF436/pptx_images"

# ══════════════════════════════════════════════════════════════════════
# PALETTE — Dark Mode 2026
# ══════════════════════════════════════════════════════════════════════
BG_DARK   = RGBColor(0x0D, 0x18, 0x2E)   # fond principal
BG_CARD   = RGBColor(0x16, 0x25, 0x40)   # carte / box
BG_CARD2  = RGBColor(0x1C, 0x2E, 0x4A)   # carte légèrement plus claire
GOLD      = RGBColor(0xF5, 0xAA, 0x00)   # accent or
TEAL      = RGBColor(0x00, 0xD4, 0xAA)   # accent teal/vert
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE  = RGBColor(0xE8, 0xEE, 0xF8)
GRAY_LT   = RGBColor(0x8A, 0x9B, 0xBB)
GREEN_OK  = RGBColor(0x2E, 0xCC, 0x71)
ORANGE_M  = RGBColor(0xE6, 0x7E, 0x22)
RED_BAD   = RGBColor(0xE7, 0x4C, 0x3C)
BLUE_ACC  = RGBColor(0x27, 0x6F, 0xBF)

# ══════════════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════════════

def bg(slide, color=BG_DARK):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def rect(slide, l, t, w, h, fill, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = lw
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size,
        bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    return tb

def img(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path,
            Inches(l), Inches(t), Inches(w), Inches(h))

def header(slide, title, subtitle=""):
    """Bande de titre sombre avec accent doré."""
    rect(slide, 0, 0, 13.33, 1.5, BG_CARD)
    rect(slide, 0, 1.5, 13.33, 0.05, GOLD)
    txt(slide, title, 0.45, 0.15, 12.4, 1.1, 30,
        bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.45, 0.9, 12.4, 0.55, 14,
            color=GOLD, italic=True)

def kpi_box(slide, l, t, w, h, big_num, label, color=GOLD, bg_c=BG_CARD):
    """Boîte KPI grande valeur + label."""
    rect(slide, l, t, w, h, bg_c)
    rect(slide, l, t, 0.15, h, color)
    txt(slide, big_num, l + 0.25, t + 0.12, w - 0.35, h * 0.5,
        size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(slide, label, l + 0.25, t + h * 0.5, w - 0.35, h * 0.45,
        size=12, color=OFFWHITE, align=PP_ALIGN.CENTER)

def bullet_card(slide, l, t, w, h, title, items, accent=GOLD):
    """Carte bento avec titre + liste de points."""
    rect(slide, l, t, w, h, BG_CARD)
    rect(slide, l, t, w, 0.48, accent)
    txt(slide, title, l + 0.15, t + 0.07, w - 0.2, 0.38,
        size=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    y = t + 0.6
    for item in items:
        rect(slide, l + 0.22, y + 0.1, 0.1, 0.1, accent)
        txt(slide, item, l + 0.42, y, w - 0.55, 0.42,
            size=12.5, color=OFFWHITE)
        y += 0.45

def progress_bar(slide, current, total=12):
    """Barre de progression discrète en bas de slide."""
    rect(slide, 0, 7.3, 13.33, 0.2, BG_CARD)
    w = 13.33 * (current / total)
    rect(slide, 0, 7.3, w, 0.2, GOLD)
    txt(slide, f"{current} / {total}", 11.8, 7.3, 1.4, 0.2,
        size=8, color=GRAY_LT, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════
# PRÉSENTATION
# ══════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ─────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITRE
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)

# Bande latérale dorée
rect(s, 0, 0, 0.4, 7.5, GOLD)
# Bande en bas
rect(s, 0, 6.8, 13.33, 0.7, BG_CARD)

# Gros titre
txt(s, "Classification du", 0.7, 1.2, 12.0, 1.0,
    size=40, bold=True, color=OFFWHITE)
txt(s, "Score de Crédit", 0.7, 2.2, 12.0, 1.1,
    size=52, bold=True, color=GOLD)

# Sous-titre
txt(s,
    "Forage de données · Machine Learning · Déploiement Web",
    0.7, 3.5, 12.0, 0.6, size=18, italic=True, color=GRAY_LT)

# Trait décoratif
rect(s, 0.7, 4.25, 5.5, 0.05, TEAL)

# Info cours / équipe
txt(s, "8INF436 — Forage de données", 0.7, 4.4, 8, 0.45,
    size=15, color=OFFWHITE)
txt(s,
    "Frantxa Cabrejos  ·  Loup-Djabril Le Bivic  ·  Nathan Razafindratsima",
    0.7, 4.9, 11, 0.45, size=13, color=GRAY_LT)
txt(s, "UQAC  —  2025", 0.7, 5.35, 5, 0.4, size=12, color=GRAY_LT)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 2 — PROBLÉMATIQUE
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
header(s, "Comment évaluer le risque de crédit de 100 000 clients ?",
       "Le problème que résout ce projet")
progress_bar(s, 2)

# Grande question visuelle
rect(s, 0.5, 1.75, 12.33, 2.5, BG_CARD)
txt(s, "Un agent humain met ~15 min par dossier.",
    0.8, 1.9, 11.5, 0.7, size=22, bold=True, color=WHITE)
txt(s, "Soit 25 000 heures pour traiter la base entière.",
    0.8, 2.5, 11.5, 0.7, size=20, italic=True, color=GOLD)
txt(s, "Notre modèle prédit les 100 000 scores en quelques secondes.",
    0.8, 3.1, 11.5, 0.7, size=18, color=TEAL)

# 3 classes en bento
for i, (label, desc, color) in enumerate([
    ("GOOD",     "Client fiable · faible risque",      GREEN_OK),
    ("STANDARD", "Profil à surveiller · risque moyen", ORANGE_M),
    ("POOR",     "Risque élevé · crédit déconseillé",  RED_BAD),
]):
    x = 0.5 + i * 4.16
    rect(s, x, 4.45, 3.9, 0.7, color)
    txt(s, label, x, 4.48, 3.9, 0.38, size=18,
        bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    txt(s, desc, x, 4.9, 3.9, 0.35, size=11.5,
        color=BG_DARK, align=PP_ALIGN.CENTER)

txt(s, "→ Classification automatique multi-classes par machine learning",
    0.5, 5.35, 12.33, 0.5, size=14, bold=True, color=OFFWHITE)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 3 — LE JEU DE DONNÉES
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
header(s, "Un dataset de 100 000 clients bancaires réels",
       "Source : Kaggle Credit Score Classification")
progress_bar(s, 3)

# 3 KPI en haut
for i, (val, lab, col) in enumerate([
    ("100 000", "observations", GOLD),
    ("28",      "variables (features)", TEAL),
    ("3",       "classes cibles", GREEN_OK),
]):
    kpi_box(s, 0.35 + i * 3.0, 1.65, 2.7, 1.1, val, lab, col)

# Graphique distribution (grande image)
img(s, f"{IMG}/img_01_cell13.png", 0.35, 2.88, 12.6, 4.15)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 4 — QUALITÉ DES DONNÉES
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
header(s, "Des données volontairement bruitées — à nettoyer",
       "Exploration de la qualité avant traitement")
progress_bar(s, 4)

# Image valeurs manquantes (grande)
img(s, f"{IMG}/img_00_cell11.png", 0.35, 1.65, 8.5, 3.7)

# 3 cartes d'info à droite
for i, (icon, title, desc) in enumerate([
    ("⚠", "Données corrompues",
     "Âges négatifs (−500), revenus avec '_', caractères spéciaux dans les champs texte"),
    ("🔧", "Nettoyage Regex",
     "Expression régulière [^0-9.-] → supprime tout sauf chiffres et décimaux"),
    ("✓", "Aucune colonne supprimée",
     "Max 15% de manquants · sous le seuil de 60% → toutes les variables conservées"),
]):
    y = 1.65 + i * 1.35
    rect(s, 9.1, y, 4.0, 1.25, BG_CARD)
    rect(s, 9.1, y, 0.14, 1.25, GOLD if i < 2 else TEAL)
    txt(s, title, 9.35, y + 0.08, 3.6, 0.38,
        size=13, bold=True, color=GOLD if i < 2 else TEAL)
    txt(s, desc, 9.35, y + 0.48, 3.6, 0.68,
        size=11.5, color=OFFWHITE, wrap=True)

txt(s, "Après nettoyage : données prêtes pour le prétraitement →",
    0.35, 5.5, 12.6, 0.45, size=13, bold=True, color=TEAL)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 5 — PIPELINE DE PRÉTRAITEMENT (visuel flux)
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
header(s, "Le pipeline complet : du brut au modèle",
       "Chaque étape est sauvegardée (.pkl) pour le déploiement")
progress_bar(s, 5)

steps = [
    ("1", "Encodage",      "LabelEncoder\n5 variables\ncatégorielles",  GOLD),
    ("2", "Imputation",    "Médiane\n→ 0 valeur\nmanquante",            TEAL),
    ("3", "Normalisation", "MinMaxScaler\n[0, 1]\npour le SVM",         BLUE_ACC),
    ("4", "Split 80/20",   "Stratifié\nproportions\nconservées",        ORANGE_M),
    ("5", "SMOTE",         "80K → 127K\n3 classes\nbalancées",          GREEN_OK),
    ("6", "Réduction",     "Feature sel.\nou PCA\nselon modèle",        RGBColor(0xAB, 0x47, 0xBC)),
]

for i, (num, title, desc, col) in enumerate(steps):
    x = 0.35 + i * 2.17
    # Boite
    rect(s, x, 1.75, 1.95, 3.8, BG_CARD)
    # En-tête coloré
    rect(s, x, 1.75, 1.95, 0.62, col)
    # Numéro
    txt(s, num, x, 1.77, 1.95, 0.55, size=22,
        bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    # Titre
    txt(s, title, x + 0.1, 2.45, 1.75, 0.45,
        size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    # Description
    txt(s, desc, x + 0.08, 2.98, 1.79, 1.5,
        size=12, color=OFFWHITE, align=PP_ALIGN.CENTER)
    # Flèche sauf dernier
    if i < 5:
        txt(s, "→", x + 1.95, 3.0, 0.22, 0.5,
            size=18, bold=True, color=GRAY_LT, align=PP_ALIGN.CENTER)

txt(s, "Pipeline complet enregistré avec pickle → utilisé directement dans l'application Streamlit",
    0.35, 5.75, 12.6, 0.5, size=13, italic=True, color=GRAY_LT)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 6 — SMOTE : RÉÉQUILIBRAGE
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
header(s, "SMOTE : on a généré 47 617 exemples synthétiques",
       "Sans équilibrage, le modèle ignorerait la classe minoritaire 'Good'")
progress_bar(s, 6)

# Graphique SMOTE
img(s, f"{IMG}/img_02_cell26.png", 0.35, 1.65, 8.5, 4.3)

# Stats à droite
kpi_box(s, 9.1, 1.65, 3.9, 1.2, "80 000", "exemples avant SMOTE", ORANGE_M)
kpi_box(s, 9.1, 2.95, 3.9, 1.2, "127 617", "exemples après SMOTE", GREEN_OK)

rect(s, 9.1, 4.25, 3.9, 1.7, BG_CARD)
rect(s, 9.1, 4.25, 0.14, 1.7, TEAL)
txt(s, "Pourquoi uniquement sur le train ?", 9.3, 4.32, 3.55, 0.42,
    size=12, bold=True, color=TEAL)
txt(s,
    "Appliquer SMOTE sur le test\nfausserait l'évaluation réelle.\nOn simule des données nouvelles\nseulement pour l'entraînement.",
    9.3, 4.78, 3.55, 1.1, size=11.5, color=OFFWHITE)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 7 — RÉDUCTION DE DIMENSIONNALITÉ
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
header(s, "Moins de dimensions, même information",
       "Deux stratégies selon le modèle utilisé")
progress_bar(s, 7)

# Colonne gauche : Feature Selection
rect(s, 0.35, 1.65, 6.1, 5.65, BG_CARD)
rect(s, 0.35, 1.65, 6.1, 0.5, GOLD)
txt(s, "SÉLECTION DE VARIABLES", 0.45, 1.68, 5.9, 0.42,
    size=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
img(s, f"{IMG}/img_03_cell29.png", 0.4, 2.22, 6.0, 2.7)
for b in [
    "Importances calculées par Forêt Aléatoire",
    "Top 8 features sur 22 retenues",
    "Supprime le bruit → meilleure généralisation",
    "Utilisé par : RF et Gradient Boosting",
]:
    # handled below
    pass

y = 5.0
for b in [
    "Importances via Forêt Aléatoire",
    "22 variables → 8 retenues",
    "Utilisé : RF + Gradient Boosting",
]:
    rect(s, 0.55, y + 0.1, 0.1, 0.1, GOLD)
    txt(s, b, 0.77, y, 5.4, 0.38, size=12, color=OFFWHITE)
    y += 0.4

# Colonne droite : PCA
rect(s, 6.85, 1.65, 6.1, 5.65, BG_CARD)
rect(s, 6.85, 1.65, 6.1, 0.5, TEAL)
txt(s, "ACP (PCA)", 6.95, 1.68, 5.9, 0.42,
    size=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
img(s, f"{IMG}/img_05_cell32.png", 6.9, 2.22, 6.0, 2.7)

y = 5.0
for b in [
    "95% de variance → 20 composantes",
    "22 dimensions réduites à 20",
    "Utilisé : SVM (LinearSVC)",
]:
    rect(s, 7.05, y + 0.1, 0.1, 0.1, TEAL)
    txt(s, b, 7.27, y, 5.4, 0.38, size=12, color=OFFWHITE)
    y += 0.4

# ─────────────────────────────────────────────────────────────────────
# SLIDE 8 — LES 3 MODÈLES
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
header(s, "Trois approches complémentaires testées",
       "Chaque modèle a ses forces selon les données")
progress_bar(s, 8)

models = [
    ("Forêt Aléatoire",
     GREEN_OK,
     "Ensemble de 100 arbres\nde décision indépendants",
     ["Chaque arbre vote → classe majoritaire",
      "Robuste sans normalisation",
      "8 features sélectionnées en entrée",
      "100 arbres · random_state=42"]),
    ("SVM — LinearSVC",
     ORANGE_M,
     "Séparateur à vaste marge\ndans l'espace des PCA",
     ["Choix linéaire : O(n) vs O(n³) pour RBF",
      "20 composantes PCA en entrée",
      "C=1.0 · max_iter=2000",
      "Calibré pour probabilités (cv=3)"]),
    ("Gradient Boosting",
     BLUE_ACC,
     "Arbres séquentiels :\nchacun corrige le précédent",
     ["Learning rate=0.1 (conservateur)",
      "Profondeur max=3 (stabilité)",
      "100 itérations d'apprentissage",
      "8 features sélectionnées"]),
]

for i, (name, col, subtitle, bullets) in enumerate(models):
    x = 0.35 + i * 4.33
    rect(s, x, 1.65, 4.0, 5.65, BG_CARD)
    rect(s, x, 1.65, 4.0, 0.65, col)
    txt(s, name, x + 0.1, 1.68, 3.8, 0.5, size=16,
        bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    txt(s, subtitle, x + 0.1, 2.38, 3.8, 0.75,
        size=12.5, italic=True, color=col, align=PP_ALIGN.CENTER)
    y = 3.25
    for b in bullets:
        rect(s, x + 0.22, y + 0.1, 0.1, 0.1, col)
        txt(s, b, x + 0.42, y, 3.35, 0.42, size=12, color=OFFWHITE)
        y += 0.48

# ─────────────────────────────────────────────────────────────────────
# SLIDE 9 — RÉSULTATS
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
header(s, "La Forêt Aléatoire remporte la comparaison — F1 = 76.6%",
       "Évaluation sur 20 000 exemples de test non vus pendant l'entraînement")
progress_bar(s, 9)

# Graphique comparaison (à gauche)
img(s, f"{IMG}/img_06_cell44.png", 0.35, 1.65, 8.0, 4.85)

# Tableau des résultats (à droite) - bento
rect(s, 8.6, 1.65, 4.5, 0.55, BG_CARD2)
for j, col_head in enumerate(["Modèle", "Acc.", "F1"]):
    x_col = [8.7, 10.6, 11.8][j]
    txt(s, col_head, x_col, 1.7, [1.7, 0.95, 0.95][j], 0.4,
        size=12, bold=True, color=GOLD,
        align=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER][j])

results = [
    ("Forêt Aléatoire",   "76.6%", "76.6%", GREEN_OK, True),
    ("Gradient Boosting", "68.5%", "68.7%", ORANGE_M, False),
    ("SVM LinearSVC",     "53.2%", "52.9%", RED_BAD,  False),
]
for i, (name, acc, f1, col, best) in enumerate(results):
    y_r = 2.3 + i * 0.82
    bg_c = RGBColor(0x1A, 0x38, 0x25) if best else BG_CARD
    rect(s, 8.6, y_r, 4.5, 0.75, bg_c)
    rect(s, 8.6, y_r, 0.2, 0.75, col)
    txt(s, name, 8.88, y_r + 0.17, 1.65, 0.42, size=12, color=OFFWHITE)
    txt(s, acc,  10.6, y_r + 0.17, 0.95, 0.42, size=13, bold=True,
        color=col, align=PP_ALIGN.CENTER)
    txt(s, f1,   11.8, y_r + 0.17, 0.95, 0.42, size=13, bold=True,
        color=col, align=PP_ALIGN.CENTER)
    if best:
        txt(s, "★ MEILLEUR", 8.65, y_r + 0.55, 1.5, 0.25,
            size=9, bold=True, color=GOLD)

# Encart CV
rect(s, 8.6, 4.83, 4.5, 1.65, BG_CARD)
rect(s, 8.6, 4.83, 0.14, 1.65, TEAL)
txt(s, "Validation croisée (k=5)", 8.82, 4.9, 4.15, 0.4,
    size=12, bold=True, color=TEAL)
txt(s,
    "RF : F1 = 82.4% ± 0.003\nGB : F1 = 75.8% ± 0.002\nSVM : F1 = 60.0% ± 0.002",
    8.82, 5.38, 4.15, 1.0, size=12.5, color=OFFWHITE)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 10 — MATRICES DE CONFUSION
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
header(s, "La Forêt Aléatoire fait le moins d'erreurs croisées",
       "Lecture : diagonale = bonnes prédictions · hors-diagonale = erreurs")
progress_bar(s, 10)

# Grande image
img(s, f"{IMG}/img_07_cell45.png", 0.35, 1.65, 12.6, 4.65)

# Ligne d'explication
rect(s, 0.35, 6.45, 12.6, 0.75, BG_CARD)
txt(s,
    "RF classe correctement 8 213 'Standard', 4 528 'Poor' et 2 577 'Good' · "
    "Le SVM confond davantage de classes → moins fiable",
    0.55, 6.52, 12.2, 0.6, size=13, italic=True, color=OFFWHITE)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 11 — APPLICATION WEB
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
header(s, "Une application web pour utiliser les modèles en production",
       "Déploiement Streamlit — 3 onglets fonctionnels")
progress_bar(s, 11)

tabs = [
    ("Onglet 1", "Saisie manuelle", GREEN_OK,
     ["Entrer les infos d'un client",
      "3 modèles prédisent en temps réel",
      "Vote majoritaire → décision finale",
      "Code couleur : vert / orange / rouge"]),
    ("Onglet 2", "Import CSV en masse", GOLD,
     ["Importer un fichier multi-clients",
      "Prédictions automatiques en batch",
      "Export CSV des résultats",
      "Idéal pour traitement industriel"]),
    ("Onglet 3", "Tableau de bord", TEAL,
     ["Accuracy · Précision · Rappel · F1",
      "Graphique comparatif interactif",
      "Identification du meilleur modèle",
      "Métriques chargées depuis le CSV"]),
]

for i, (tag, title, col, bullets) in enumerate(tabs):
    x = 0.35 + i * 4.33
    rect(s, x, 1.65, 4.0, 5.65, BG_CARD)
    rect(s, x, 1.65, 4.0, 0.55, col)
    txt(s, tag,   x + 0.1, 1.67, 3.8, 0.28, size=10, color=BG_DARK,
        align=PP_ALIGN.CENTER)
    txt(s, title, x + 0.1, 1.9, 3.8, 0.38, size=14, bold=True,
        color=BG_DARK, align=PP_ALIGN.CENTER)
    y = 2.32
    for b in bullets:
        rect(s, x + 0.22, y + 0.1, 0.1, 0.1, col)
        txt(s, b, x + 0.42, y, 3.35, 0.42, size=12.5, color=OFFWHITE)
        y += 0.5

rect(s, 0.35, 6.15, 12.63, 0.55, RGBColor(0x12, 0x2D, 0x48))
rect(s, 0.35, 6.15, 0.2,   0.55, BLUE_ACC)
txt(s,
    "Modèles chargés via pickle · pipeline identique à l'entraînement · zéro ré-entraînement nécessaire",
    0.65, 6.2, 12.0, 0.45, size=12.5, italic=True, color=OFFWHITE)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 12 — CONCLUSION
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank); bg(s)
progress_bar(s, 12)

# Bande dorée gauche
rect(s, 0, 0, 0.4, 7.5, GOLD)
rect(s, 0.4, 6.6, 12.93, 0.7, BG_CARD)

txt(s, "Ce qu'on retient", 0.65, 0.3, 10, 0.75,
    size=36, bold=True, color=WHITE)
txt(s, "Conclusion", 0.65, 1.05, 8, 0.5,
    size=20, italic=True, color=GOLD)
rect(s, 0.65, 1.6, 5.5, 0.05, TEAL)

bilan = [
    (GOLD,     "Pipeline complet",
     "Nettoyage → encodage → normalisation → SMOTE → réduction → modèle → déploiement"),
    (GREEN_OK, "Forêt Aléatoire = meilleur modèle",
     "F1 pondéré de 76.6% sur test · 82.4% en validation croisée (k=5)"),
    (TEAL,     "SMOTE a été déterminant",
     "Sans lui, les classes minoritaires auraient été ignorées par les algorithmes"),
    (ORANGE_M, "Application prête à l'emploi",
     "Streamlit déployable · prédictions individuelles et en masse · interface intuitive"),
]

for i, (col, title, desc) in enumerate(bilan):
    y = 1.85 + i * 1.1
    rect(s, 0.65, y, 0.5, 0.5, col)
    txt(s, str(i + 1), 0.65, y + 0.03, 0.5, 0.44, size=20, bold=True,
        color=BG_DARK, align=PP_ALIGN.CENTER)
    txt(s, title, 1.3, y + 0.02, 11.5, 0.38, size=14,
        bold=True, color=col)
    txt(s, desc,  1.3, y + 0.48, 11.5, 0.52, size=12.5, color=OFFWHITE)

txt(s, "Merci pour votre attention  ·  Des questions ?",
    0.4, 6.63, 12.93, 0.55, size=20, bold=True, color=GOLD,
    align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════
# SAUVEGARDE
# ══════════════════════════════════════════════════════════════════════
out = "/home/user/projet-forage-8INF436/presentation_credit_score_v2.pptx"
prs.save(out)
print(f"✓ Fichier créé : {out}")
print(f"✓ Slides : {len(prs.slides)}")
print(f"✓ Taille : {os.path.getsize(out) / 1024:.0f} Ko")
